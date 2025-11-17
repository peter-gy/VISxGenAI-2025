import marimo

__generated_with = "0.17.8"
app = marimo.App(width="columns")


@app.cell(column=0, hide_code=True)
def _(mo):
    mo.image(agent_diagram_url("01"))
    return


@app.cell
def _(df):
    df
    return


@app.cell
def _(read_dataset):
    dataset_uri = "https://data.abudhabi/opendata/sites/default/files/uploaded_resources/student_enrollment_grade_2025.xlsx"
    df = read_dataset(dataset_uri)
    return (df,)


@app.cell
def _(lm_registry_from_env, load_env, make_session_id):
    env = load_env()
    lm_registry = lm_registry_from_env(env)
    session_id = make_session_id()
    return env, lm_registry, session_id


@app.function(hide_code=True)
def agent_diagram_url(id: str):
    return f"https://peter-gy.github.io/VISxGenAI-2025/assets/agents/{id}.svg"


@app.cell
def _():
    import marimo as mo
    from visxgenai.core import (
        lm_registry_from_env,
        load_env,
        make_session_id,
        read_dataset,
    )

    return lm_registry_from_env, load_env, make_session_id, mo, read_dataset


@app.cell(column=1, hide_code=True)
def _(mo):
    mo.image(agent_diagram_url("02"))
    return


@app.cell
def _(field_refiner_output):
    dict(field_refiner_output)
    return


@app.cell(hide_code=True)
def _(df, field_refiner_agent):
    field_refiner_output = field_refiner_agent(df=df)
    return (field_refiner_output,)


@app.cell
def _(FieldRefinerAgent, lm_registry):
    field_refiner_agent = FieldRefinerAgent(
        lm_registry.lm_initializer("vertex_ai/gemini-2.5-flash"),
        langfuse=lm_registry.langfuse,
    )
    return (field_refiner_agent,)


@app.cell
def _():
    from visxgenai.agents.field_refiner.agent import FieldRefinerAgent

    return (FieldRefinerAgent,)


@app.cell(column=2, hide_code=True)
def _(mo):
    mo.image(agent_diagram_url("03"))
    return


@app.cell
def _(dataset_describer_output):
    dict(dataset_describer_output)
    return


@app.cell
def _(dataset_describer, field_refiner_output):
    dataset_describer_output = dataset_describer(
        df=field_refiner_output["df"],
        semantic_schema=field_refiner_output["semantic_schema"],
    )
    return (dataset_describer_output,)


@app.cell
def _(DatasetDescriberAgent, lm_registry):
    dataset_describer = DatasetDescriberAgent(
        lm_registry.lm_initializer(
            "vertex_ai/gemini-2.0-flash",
        ),
        langfuse=lm_registry.langfuse,
    )
    return (dataset_describer,)


@app.cell
def _():
    from visxgenai.agents.dataset_describer.agent import DatasetDescriberAgent

    return (DatasetDescriberAgent,)


@app.cell(column=3, hide_code=True)
def _(mo):
    mo.image(agent_diagram_url("04"))
    return


@app.cell
def _(field_expander_output):
    dict(field_expander_output)
    return


@app.cell
def _(dataset_describer_output, field_expander, field_refiner_output):
    field_expander_output = field_expander(
        df=field_refiner_output["df"],
        dataset_summary=str(dict(dataset_describer_output)),
        dataset_fields=dataset_describer_output["fields"],
        field_info=field_refiner_output["refinements"],
    )
    return (field_expander_output,)


@app.cell
def _(FieldExpanderAgent, lm_registry):
    field_expander = FieldExpanderAgent(
        lm_registry.lm_initializer(
            "gemini/gemini-2.0-flash",
            max_tokens=8_000,
            tools=[
                {"googleSearch": {}},
            ],
        ),
        langfuse=lm_registry.langfuse,
    )
    return (field_expander,)


@app.cell
def _():
    from visxgenai.agents.field_expander.agent import FieldExpanderAgent

    return (FieldExpanderAgent,)


@app.cell(column=4, hide_code=True)
def _(mo):
    mo.image(agent_diagram_url("05"))
    return


@app.cell
def _(dataset_profiler_output):
    dict(dataset_profiler_output)
    return


@app.cell
def _(dataset_profiler, field_expander_output, field_refiner_output):
    dataset_profiler_output = dataset_profiler(
        df=field_expander_output["df"],
        semantic_schema=field_refiner_output["semantic_schema"],
    )
    return (dataset_profiler_output,)


@app.cell
def _(DatasetProfilerAgent, lm_registry):
    dataset_profiler = DatasetProfilerAgent(
        langfuse=lm_registry.langfuse,
    )
    return (dataset_profiler,)


@app.cell
def _():
    from visxgenai.agents.dataset_profiler.agent import DatasetProfilerAgent

    return (DatasetProfilerAgent,)


@app.cell(column=5, hide_code=True)
def _(mo):
    mo.image(agent_diagram_url("06"))
    return


@app.cell
def _(insight_planner_output):
    [m.model_dump() for m in insight_planner_output.plan]
    return


@app.cell
def _(InsightPlanTools, insight_planner_output):
    insight_plan_tools = InsightPlanTools(insight_planner_output.plan)
    return (insight_plan_tools,)


@app.cell
def _(dataset_describer_output, dataset_profiler_output, insight_planner):
    insight_planner_output = insight_planner(
        dataset_title=dataset_describer_output.title,
        dataset_description=dataset_describer_output.summary,
        field_descriptions=dataset_describer_output.fields,
        dataset_profile=dataset_profiler_output.profile,
    )
    insight_planner_output
    return (insight_planner_output,)


@app.cell
def _(InsightPlannerAgent, lm_registry):
    insight_planner = InsightPlannerAgent(
        init_planner_lm=lm_registry.lm_initializer(
            "vertex_ai/gemini-2.5-flash",
            max_tokens=32_000,
            # reasoning_effort="disable",
        ),
        init_judge_lm=lm_registry.lm_initializer(
            "vertex_ai/gemini-2.5-flash-lite",
            max_tokens=16_000,
            reasoning_effort="disable",
        ),
        langfuse=lm_registry.langfuse,
    )
    return (insight_planner,)


@app.cell
def _(
    construct_dataset_context,
    dataset_describer_output,
    dataset_profiler_output,
):
    construct_dataset_context(
        dataset_title=dataset_describer_output.title,
        dataset_description=dataset_describer_output.summary,
        field_descriptions=dataset_describer_output.fields,
        dataset_profile=dataset_profiler_output.profile,
    )
    return


@app.cell
def _():
    from visxgenai.agents.insight_planner.utils import construct_dataset_context

    return (construct_dataset_context,)


@app.cell
def _():
    from visxgenai.agents.insight_planner.agent import InsightPlannerAgent
    from visxgenai.agents.insight_planner.tools import InsightPlanTools

    return InsightPlanTools, InsightPlannerAgent


@app.cell(column=6, hide_code=True)
def _(mo):
    mo.image(agent_diagram_url("07"))
    return


@app.cell
def _(datset_deriver_output):
    materialized_insights = [
        {"goal": ds["goal"]} | dict(zip(("dataset", "sql"), ds["materialize"]()))
        for ds in datset_deriver_output.datasets
    ]
    materialized_insights
    return (materialized_insights,)


@app.cell
def _(datset_deriver_output):
    dict(datset_deriver_output)
    return


@app.cell
def _(datset_deriver_output):
    datset_deriver_output.queries
    return


@app.cell
def _(datset_deriver_output):
    dict(datset_deriver_output)
    return


@app.cell
def _(
    dataset_describer_output,
    dataset_profiler_output,
    datset_deriver,
    field_expander_output,
    insight_plan_tools,
):
    datset_deriver_output = datset_deriver(
        df=field_expander_output.df,
        dataset_title=dataset_describer_output.title,
        dataset_description=dataset_describer_output.summary,
        field_descriptions=dataset_describer_output.fields,
        dataset_profile=dataset_profiler_output.profile,
        nl_queries=insight_plan_tools.list_goals(),
    )
    datset_deriver_output
    return (datset_deriver_output,)


@app.cell
def _(DatasetDeriverAgent, lm_registry):
    datset_deriver = DatasetDeriverAgent(
        init_drafter_lm=lm_registry.lm_initializer(
            "vertex_ai/gemini-2.5-pro",
            max_tokens=60_000,
        ),
        init_repairer_lm=lm_registry.lm_initializer(
            "vertex_ai/gemini-2.5-flash",
            max_tokens=60_000,
        ),
        langfuse=lm_registry.langfuse,
    )
    return (datset_deriver,)


@app.cell
def _():
    from visxgenai.agents.dataset_deriver.agent import DatasetDeriverAgent

    return (DatasetDeriverAgent,)


@app.cell(column=7, hide_code=True)
def _(mo):
    mo.image(agent_diagram_url("08"))
    return


@app.cell
def _(dataset_publisher_output):
    dict(dataset_publisher_output)
    return


@app.cell
def _(dataset_publisher, materialized_insights, session_id):
    dataset_publisher_output = dataset_publisher(
        materialized_insights=materialized_insights,
        session=session_id,
    )
    return (dataset_publisher_output,)


@app.cell
def _(DatasetPublisherAgent, env, lm_registry):
    dataset_publisher = DatasetPublisherAgent(
        aws_access_key_id=env["AWS_ACCESS_KEY_ID"],
        aws_secret_access_key=env["AWS_SECRET_ACCESS_KEY"],
        aws_endpoint_url=env["AWS_ENDPOINT_URL"],
        cdn_base_url=env["AWS_CDN_BASE_URL"],
        bucket=env["AWS_BUCKET"],
        langfuse=lm_registry.langfuse,
    )
    return (dataset_publisher,)


@app.cell
def _():
    from visxgenai.agents.dataset_publisher.agent import DatasetPublisherAgent

    return (DatasetPublisherAgent,)


@app.cell(column=8, hide_code=True)
def _(mo):
    mo.image(agent_diagram_url("09"))
    return


@app.cell
def _(dataset_visualizer_output):
    dict(dataset_visualizer_output)
    return


@app.cell
def _(
    dataset_visualizer,
    datset_deriver_output,
    insight_plan_tools,
    materialized_insights,
):
    dataset_visualizer_output = dataset_visualizer(
        queried_insights=datset_deriver_output.datasets,
        materialized_insights=materialized_insights,
        insight_plan_tools=insight_plan_tools,
    )
    return (dataset_visualizer_output,)


@app.cell
def _(DatasetVisualizerAgent, lm_registry):
    dataset_visualizer = DatasetVisualizerAgent(
        langfuse=lm_registry.langfuse,
    )
    return (dataset_visualizer,)


@app.cell
def _():
    from visxgenai.agents.dataset_visualizer.agent import DatasetVisualizerAgent

    return (DatasetVisualizerAgent,)


@app.cell(column=9, hide_code=True)
def _(mo):
    mo.image(agent_diagram_url("10"))
    return


@app.cell
def _(report_narrator_output):
    report_narrator_output.content
    return


@app.cell
def _(
    dataset_describer_output,
    dataset_visualizer_output,
    datset_deriver_output,
    insight_plan_tools,
    report_narrator,
):
    report_narrator_output = report_narrator(
        parent_dataset_description=dataset_describer_output.summary,
        queried_insights=datset_deriver_output.datasets,
        visualized_insights=dataset_visualizer_output.recommendations,
        insight_plan_tools=insight_plan_tools,
    )
    return (report_narrator_output,)


@app.cell
def _(DataReportNarratorAgent, lm_registry):
    report_narrator = DataReportNarratorAgent(
        init_narrator_lm=lm_registry.lm_initializer(
            "vertex_ai/gemini-2.5-pro",
            max_tokens=32_000,
        ),
        init_structurer_lm=lm_registry.lm_initializer(
            "vertex_ai/gemini-2.5-flash",
            max_tokens=32_000,
        ),
        langfuse=lm_registry.langfuse,
    )
    return (report_narrator,)


@app.cell
def _():
    from visxgenai.agents.dataset_narrator.agent import DataReportNarratorAgent

    return (DataReportNarratorAgent,)


@app.cell(column=10, hide_code=True)
def _(mo):
    mo.image(agent_diagram_url("11"))
    return


@app.cell
def _(dataset_reporter_output):
    dict(dataset_reporter_output)
    return


@app.cell
def _(dataset_reporter_output, mo):
    mo.ui.code_editor(dataset_reporter_output.nb, language="html")
    return


@app.cell
def _(
    dataset_publisher_output,
    dataset_reporter,
    dataset_visualizer_output,
    datset_deriver_output,
    materialized_insights,
    report_narrator_output,
    session_id,
):
    dataset_reporter_output = dataset_reporter(
        session=session_id,
        queried_insights=datset_deriver_output.datasets,
        materialized_insights=materialized_insights,
        visualized_insights=dataset_visualizer_output.recommendations,
        published_datasets=dataset_publisher_output.records,
        report_structure=report_narrator_output.content,
        traces={"flow": "#", "observations": {}},
    )
    return (dataset_reporter_output,)


@app.cell
def _(DatasetReporterAgent, env, lm_registry):
    dataset_reporter = DatasetReporterAgent(
        nb_builder_api_username=env["NB_BUILDER_API_USERNAME"],
        nb_builder_api_password=env["NB_BUILDER_API_PASSWORD"],
        aws_access_key_id=env["AWS_ACCESS_KEY_ID"],
        aws_secret_access_key=env["AWS_SECRET_ACCESS_KEY"],
        aws_endpoint_url=env["AWS_ENDPOINT_URL"],
        aws_bucket=env["AWS_BUCKET"],
        aws_cdn_base_url=env["AWS_CDN_BASE_URL"],
        langfuse=lm_registry.langfuse,
        api_base_url=env["NB_BUILDER_API_BASE_URL"],
    )
    return (dataset_reporter,)


@app.cell
def _():
    from visxgenai.agents.dataset_reporter.agent import DatasetReporterAgent

    return (DatasetReporterAgent,)


@app.cell(column=11, hide_code=True)
def _(dataset_reporter_output, mo):
    mo.iframe(dataset_reporter_output.html)
    return


@app.cell(column=12)
def _(presentation_narrator_output):
    presentation_narrator_output.presentation_content
    return


@app.cell
def _(presentation_narrator, report_narrator_output):
    presentation_narrator_output = presentation_narrator(
        report_content=report_narrator_output.content,
    )
    return (presentation_narrator_output,)


@app.cell
def _(PresentationNarratorAgent, lm_registry):
    presentation_narrator = PresentationNarratorAgent(
        init_lm=lm_registry.lm_initializer(
            "vertex_ai/gemini-2.5-pro",
            max_tokens=32_000,
        ),
        langfuse=lm_registry.langfuse,
    )
    return (presentation_narrator,)


@app.cell
def _():
    from visxgenai.agents.dataset_narrator.agent import PresentationNarratorAgent

    return (PresentationNarratorAgent,)


@app.cell(column=13)
def _():
    from visxgenai.slide_kit import load_slide_template

    return (load_slide_template,)


@app.cell
def _(
    dataset_visualizer_output,
    presentation_narrator_output,
    render_presentation,
):
    render_presentation(
        content=presentation_narrator_output.presentation_content,
        charts_by_goal={
            rec["goal"]: rec["chart"]
            for rec in dataset_visualizer_output.recommendations
        },
    ).save("dynademo.pptx")
    return


@app.cell
def _(
    load_slide_template,
    render_conclusion_slide,
    render_insight_slide,
    render_report_intro_slide,
    render_section_intro_slide,
    render_summary_slide,
    render_title_slide,
):
    def render_presentation(
        content: dict,
        charts_by_goal: dict,
        web_report_url: str = "https://peter.gy",
        template: str = "mckinsey",
    ):
        import datetime as dt

        prs = load_slide_template("mckinsey")
        render_title_slide(
            prs,
            title=content["title"],
            subtitle=content["subtitle"],
            topic=content["domain"],
            date=dt.datetime.now().strftime("%m/%d/%Y"),
        )
        render_report_intro_slide(
            prs,
            title="Significance",
            intro=content["intro"],
        )

        for mainsection in content["sections"]:
            render_section_intro_slide(
                prs,
                title=mainsection["title"],
                intro=mainsection["intro"],
            )
            for subsection in mainsection["sections"]:
                render_insight_slide(
                    prs,
                    title=subsection["title"],
                    bulletpoints=subsection["content"],
                    chart=charts_by_goal[subsection["goal"]],
                )

        render_summary_slide(
            prs,
            key_takeaways=content["key_takeaways"],
        )
        render_conclusion_slide(
            prs,
            call_to_action_url=web_report_url,
            qrcode_image_path=make_qr_code(web_report_url),
        )

        return prs

    return (render_presentation,)


@app.cell
def _(ELEMENT_MAPPING):
    def _get_placeholder(slide, index: int):
        """Finds a placeholder by the name we set in the Selection Pane."""
        return list(slide.placeholders)[index]

    def _prepare_slide_for_render(prs, layout_name: str) -> tuple:
        layout = prs.slide_layouts.get_by_name(layout_name)
        slide = prs.slides.add_slide(layout)
        placeholder_map = ELEMENT_MAPPING[layout_name]
        phs = {
            name: _get_placeholder(slide, ph_name)
            for name, ph_name in placeholder_map.items()
        }
        return slide, phs

    def _render_bullets(ph, items: list[str]):
        if not items:
            return

        text_frame = ph.text_frame
        text_frame.clear()

        for i, item_text in enumerate(items):
            # Use the existing paragraph for the first item, add new ones for the rest
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            # Apply all text and formatting
            p.text = f"❯ {item_text}"
            p.level = 0

    def render_chart_and_replace_placeholder(chart, slide, placeholder):
        """
        Renders any Altair chart and replaces a placeholder with the result.

        This method ensures:
        1.  High-quality rendering by using a high DPI.
        2.  The entire chart is visible, fitting perfectly within the placeholder bounds.
        3.  Non-faceted charts have a standard aspect ratio to prevent sparsity.
        4.  The original placeholder (and its border) is completely removed.

        Args:
            chart (alt.Chart): The Altair chart object (simple or faceted).
            slide (pptx.slide.Slide): The slide object containing the placeholder.
            placeholder (pptx.shapes.placeholder._BasePlaceholder): The placeholder to be replaced.
            dpi (int): The resolution for rendering the chart.
        """
        import io
        import altair as alt
        from PIL import Image

        # --- Step 1: Control the shape for non-faceted charts to prevent sparsity ---
        if not isinstance(chart, alt.FacetChart):
            # Apply a standard 4:3 aspect ratio for a more "normal" look
            chart = chart.properties(width=400, height=450)

        # Increase axis label font size for better readability
        chart = chart.configure_axis(
            labelFontSize=12,
            titleFontSize=14,
        )

        # --- Step 2: Render the prepared chart at high quality ---
        image_stream = io.BytesIO()
        chart.save(image_stream, format="png", scale_factor=3)
        image_stream.seek(0)

        # --- Step 3: Calculate the best fit for the rendered image ---
        ph_left, ph_top, ph_width, ph_height = (
            placeholder.left,
            placeholder.top,
            placeholder.width,
            placeholder.height,
        )
        with Image.open(image_stream) as img:
            img_width, img_height = img.size

        aspect_ratio_img = img_width / img_height
        aspect_ratio_ph = ph_width / ph_height

        if aspect_ratio_img > aspect_ratio_ph:
            new_width = ph_width
            new_height = int(new_width / aspect_ratio_img)
        else:
            new_height = ph_height
            new_width = int(new_height * aspect_ratio_img)

        # Center the image within the original placeholder's area
        new_left = ph_left + int((ph_width - new_width) / 2)
        new_top = ph_top + int((ph_height - new_height) / 2)

        # --- Step 4: Delete the original placeholder ---
        # Access the shape's XML element and its parent, then remove it.
        sp = placeholder.element
        sp.getparent().remove(sp)

        # --- Step 5: Add the perfectly sized image to the slide ---
        slide.shapes.add_picture(
            image_stream, new_left, new_top, width=new_width, height=new_height
        )

    def render_title_slide(prs, title: str, subtitle: str, topic: str, date: str):
        slide, phs = _prepare_slide_for_render(prs, "Title Layout")

        phs["Title"].text = title
        phs["Subtitle"].text = subtitle
        phs["Topic"].text = topic
        phs["Date"].text = date

        return prs

    def render_report_intro_slide(prs, title: str, intro: str):
        slide, phs = _prepare_slide_for_render(prs, "Title and Content Layout")

        phs["Title"].text = title
        phs["Content"].text = intro

        return prs

    def render_section_intro_slide(prs, title: str, intro: str):
        slide, phs = _prepare_slide_for_render(prs, "Divider Layout")

        phs["Title"].text = title
        phs["Intro"].text = intro

        return prs

    def render_insight_slide(
        prs,
        title: str,
        bulletpoints: list[str],
        chart,
    ):
        slide, phs = _prepare_slide_for_render(prs, "Text and Chart Layout")

        phs["Title"].text = title
        chart = chart.configure_mark(color="#162645")
        render_chart_and_replace_placeholder(chart, slide, phs["Chart"])
        _render_bullets(phs["Description"], bulletpoints)

        return prs

    def render_summary_slide(prs, key_takeaways: list[str]):
        slide, phs = _prepare_slide_for_render(prs, "Summary Layout")

        _render_bullets(phs["Summary"], key_takeaways)

        return prs

    def render_conclusion_slide(
        prs,
        call_to_action_url: str,
        qrcode_image_path: str,
        text: str = "Interactive Report Available Here",
    ):
        from pptx.dml.color import RGBColor

        slide, phs = _prepare_slide_for_render(prs, "Explore More Layout")

        phs["CallToAction"].click_action.hyperlink.address = call_to_action_url
        phs["CallToAction"].text = text
        phs["CallToAction"].text_frame.paragraphs[0].font.underline = True

        phs["QR"].insert_picture(qrcode_image_path)

        return prs

    return (
        render_conclusion_slide,
        render_insight_slide,
        render_report_intro_slide,
        render_section_intro_slide,
        render_summary_slide,
        render_title_slide,
    )


@app.function
def make_qr_code(data: str) -> str:
    import qrcode
    import io

    # Create QR code with custom settings for better quality
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=64,
        border=4,
    )

    # Add data
    qr.add_data(data)
    qr.make(fit=True)

    # Create image with custom colors
    img = qr.make_image(fill_color="#162645", back_color="white")

    # Create stream
    image_stream = io.BytesIO()
    img.save(image_stream, format="PNG")
    image_stream.seek(0)

    return image_stream


@app.cell
def _():
    ELEMENT_MAPPING = {
        "Title Layout": {
            "Title": 0,
            "Subtitle": 1,
            "Topic": 2,
            "Date": 3,
        },
        "Title and Content Layout": {
            "Title": 0,
            "Content": 1,
        },
        "Divider Layout": {
            "Title": 0,
            "Intro": 1,
        },
        "Text and Chart Layout": {
            "Title": 0,
            "Description": 1,
            "Chart": 2,
        },
        "Summary Layout": {
            "Summary": 0,
        },
        "Explore More Layout": {
            "CallToAction": 0,
            "QR": 1,
        },
    }

    return (ELEMENT_MAPPING,)


if __name__ == "__main__":
    app.run()
