from pptx import Presentation

from .templates import load_slide_template

ELEMENT_MAPPING = {
    "Title Layout": {
        "Title": 0,
        "Subtitle": 1,
        "Topic": 2,
        "Date": 3,
    },
    "Title and Content Layout": {
        "Title": 0,
        "Content": 1,
    },
    "Divider Layout": {
        "Title": 0,
        "Intro": 1,
    },
    "Text and Chart Layout": {
        "Title": 0,
        "Description": 1,
        "Chart": 2,
    },
    "Summary Layout": {
        "Summary": 0,
    },
    "Explore More Layout": {
        "CallToAction": 0,
        "QR": 1,
    },
}


def make_qr_code(data: str) -> str:
    import io

    import qrcode

    # Create QR code with custom settings for better quality
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=64,
        border=4,
    )

    # Add data
    qr.add_data(data)
    qr.make(fit=True)

    # Create image with custom colors
    img = qr.make_image(fill_color="#162645", back_color="white")

    # Create stream
    image_stream = io.BytesIO()
    img.save(image_stream, format="PNG")
    image_stream.seek(0)

    return image_stream


def _get_placeholder(slide, index: int):
    """Finds a placeholder by the name we set in the Selection Pane."""
    return list(slide.placeholders)[index]


def _prepare_slide_for_render(prs, layout_name: str) -> tuple:
    layout = prs.slide_layouts.get_by_name(layout_name)
    slide = prs.slides.add_slide(layout)
    placeholder_map = ELEMENT_MAPPING[layout_name]
    phs = {
        name: _get_placeholder(slide, ph_name)
        for name, ph_name in placeholder_map.items()
    }
    return slide, phs


def _render_bullets(ph, items: list[str]):
    if not items:
        return

    text_frame = ph.text_frame
    text_frame.clear()

    for i, item_text in enumerate(items):
        # Use the existing paragraph for the first item, add new ones for the rest
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # Apply all text and formatting
        p.text = f"❯ {item_text}"
        p.level = 0


def render_chart_and_replace_placeholder(chart, slide, placeholder):
    """
    Renders any Altair chart and replaces a placeholder with the result.

    This method ensures:
    1.  High-quality rendering by using a high DPI.
    2.  The entire chart is visible, fitting perfectly within the placeholder bounds.
    3.  Non-faceted charts have a standard aspect ratio to prevent sparsity.
    4.  The original placeholder (and its border) is completely removed.

    Args:
        chart (alt.Chart): The Altair chart object (simple or faceted).
        slide (pptx.slide.Slide): The slide object containing the placeholder.
        placeholder (pptx.shapes.placeholder._BasePlaceholder): The placeholder to be replaced.
        dpi (int): The resolution for rendering the chart.
    """
    import io

    import altair as alt
    from PIL import Image

    # --- Step 1: Control the shape for non-faceted charts to prevent sparsity ---
    if not isinstance(chart, alt.FacetChart):
        # Apply a standard 4:3 aspect ratio for a more "normal" look
        chart = chart.properties(width=400, height=450)

    # Increase axis label font size for better readability
    chart = chart.configure_axis(
        labelFontSize=12,
        titleFontSize=14,
    )

    # --- Step 2: Render the prepared chart at high quality ---
    image_stream = io.BytesIO()
    chart.save(image_stream, format="png", scale_factor=3)
    image_stream.seek(0)

    # --- Step 3: Calculate the best fit for the rendered image ---
    ph_left, ph_top, ph_width, ph_height = (
        placeholder.left,
        placeholder.top,
        placeholder.width,
        placeholder.height,
    )
    with Image.open(image_stream) as img:
        img_width, img_height = img.size

    aspect_ratio_img = img_width / img_height
    aspect_ratio_ph = ph_width / ph_height

    if aspect_ratio_img > aspect_ratio_ph:
        new_width = ph_width
        new_height = int(new_width / aspect_ratio_img)
    else:
        new_height = ph_height
        new_width = int(new_height * aspect_ratio_img)

    # Center the image within the original placeholder's area
    new_left = ph_left + int((ph_width - new_width) / 2)
    new_top = ph_top + int((ph_height - new_height) / 2)

    # --- Step 4: Delete the original placeholder ---
    # Access the shape's XML element and its parent, then remove it.
    sp = placeholder.element
    sp.getparent().remove(sp)

    # --- Step 5: Add the perfectly sized image to the slide ---
    slide.shapes.add_picture(
        image_stream, new_left, new_top, width=new_width, height=new_height
    )


def render_title_slide(prs, title: str, subtitle: str, topic: str, date: str):
    slide, phs = _prepare_slide_for_render(prs, "Title Layout")

    phs["Title"].text = title
    phs["Subtitle"].text = subtitle
    phs["Topic"].text = topic
    phs["Date"].text = date

    return prs


def render_report_intro_slide(prs, title: str, intro: str):
    slide, phs = _prepare_slide_for_render(prs, "Title and Content Layout")

    phs["Title"].text = title
    phs["Content"].text = intro

    return prs


def render_section_intro_slide(prs, title: str, intro: str):
    slide, phs = _prepare_slide_for_render(prs, "Divider Layout")

    phs["Title"].text = title
    phs["Intro"].text = intro

    return prs


def render_insight_slide(
    prs,
    title: str,
    bulletpoints: list[str],
    chart,
):
    slide, phs = _prepare_slide_for_render(prs, "Text and Chart Layout")

    phs["Title"].text = title
    chart = chart.configure_mark(color="#162645")
    render_chart_and_replace_placeholder(chart, slide, phs["Chart"])
    _render_bullets(phs["Description"], bulletpoints)

    return prs


def render_summary_slide(prs, key_takeaways: list[str]):
    slide, phs = _prepare_slide_for_render(prs, "Summary Layout")

    _render_bullets(phs["Summary"], key_takeaways)

    return prs


def render_conclusion_slide(
    prs,
    call_to_action_url: str,
    qrcode_image_path: str,
    text: str = "Interactive Report Available Here",
):
    slide, phs = _prepare_slide_for_render(prs, "Explore More Layout")

    phs["CallToAction"].click_action.hyperlink.address = call_to_action_url
    phs["CallToAction"].text = text
    phs["CallToAction"].text_frame.paragraphs[0].font.underline = True

    phs["QR"].insert_picture(qrcode_image_path)

    return prs


def render_presentation(
    content: dict,
    charts_by_goal: dict,
    web_report_url: str = "https://peter.gy",
    template: str = "mckinsey",
) -> Presentation:
    import datetime as dt

    prs = load_slide_template(template)
    render_title_slide(
        prs,
        title=content["title"],
        subtitle=content["subtitle"],
        topic=content["domain"],
        date=dt.datetime.now().strftime("%m/%d/%Y"),
    )
    render_report_intro_slide(
        prs,
        title="Significance",
        intro=content["intro"],
    )

    for mainsection in content["sections"]:
        render_section_intro_slide(
            prs,
            title=mainsection["title"],
            intro=mainsection["intro"],
        )
        for subsection in mainsection["sections"]:
            render_insight_slide(
                prs,
                title=subsection["title"],
                bulletpoints=subsection["content"],
                chart=charts_by_goal[subsection["goal"]],
            )

    render_summary_slide(
        prs,
        key_takeaways=content["key_takeaways"],
    )
    render_conclusion_slide(
        prs,
        call_to_action_url=web_report_url,
        qrcode_image_path=make_qr_code(web_report_url),
    )

    return prs


__all__ = ["render_presentation"]
